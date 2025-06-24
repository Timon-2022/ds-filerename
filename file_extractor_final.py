import os
import logging
from pathlib import Path
from typing import Optional, Dict, Any
import chardet

# 文档处理库
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False

try:
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)

class FileContentExtractor:
    """文件内容提取器，支持多种文件格式"""
    
    # 支持的文件类型
    SUPPORTED_EXTENSIONS = {
        '.txt': 'text',
        '.md': 'markdown',
        '.py': 'code',
        '.js': 'code',
        '.html': 'code',
        '.css': 'code',
        '.json': 'code',
        '.xml': 'code',
        '.csv': 'text',
        '.log': 'text',
        '.ini': 'text',
        '.cfg': 'text',
        '.conf': 'text',
        '.yml': 'text',
        '.yaml': 'text',
        '.sql': 'code',
        '.sh': 'code',
        '.bat': 'code',
        '.ps1': 'code',
        '.docx': 'document',
        '.pdf': 'document',
        '.xlsx': 'spreadsheet',
        '.xls': 'spreadsheet',
        '.pptx': 'presentation',
        '.ppt': 'presentation'
    }
    
    def __init__(self):
        self.max_content_length = 5000  # 最大内容长度
    
    def is_supported_file(self, file_path: Path) -> bool:
        """检查文件是否支持内容提取"""
        return file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS
    
    def get_file_type(self, file_path: Path) -> str:
        """获取文件类型"""
        ext = file_path.suffix.lower()
        return self.SUPPORTED_EXTENSIONS.get(ext, 'unknown')
    
    def detect_encoding(self, file_path: Path) -> str:
        """检测文件编码"""
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)  # 读取前10KB用于检测编码
                result = chardet.detect(raw_data)
                return result.get('encoding', 'utf-8') or 'utf-8'
        except Exception:
            return 'utf-8'
    
    def extract_text_content(self, file_path: Path) -> Optional[str]:
        """提取纯文本文件内容"""
        try:
            encoding = self.detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                content = f.read(self.max_content_length)
                return content
        except Exception as e:
            logger.error(f"提取文本文件内容失败 {file_path}: {str(e)}")
            return None
    
    def extract_markdown_content(self, file_path: Path) -> Optional[str]:
        """提取Markdown文件内容"""
        try:
            if not MARKDOWN_AVAILABLE:
                return self.extract_text_content(file_path)
            
            encoding = self.detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                md_content = f.read(self.max_content_length)
                # 简单处理：移除Markdown标记，保留纯文本
                import re
                # 移除图片链接
                md_content = re.sub(r'!\[.*?\]\(.*?\)', '', md_content)
                # 移除链接
                md_content = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', md_content)
                # 移除标题标记
                md_content = re.sub(r'^#+\s*', '', md_content, flags=re.MULTILINE)
                # 移除代码块标记
                md_content = re.sub(r'```.*?```', '', md_content, flags=re.DOTALL)
                md_content = re.sub(r'`([^`]+)`', r'\1', md_content)
                # 移除多余空行
                md_content = re.sub(r'\n\s*\n', '\n\n', md_content)
                
                return md_content.strip()
        except Exception as e:
            logger.error(f"提取Markdown文件内容失败 {file_path}: {str(e)}")
            return None
    
    def extract_docx_content(self, file_path: Path) -> Optional[str]:
        """提取Word文档内容"""
        try:
            if not DOCX_AVAILABLE:
                logger.warning("python-docx 未安装，无法处理 .docx 文件")
                return None
            
            doc = Document(file_path)
            content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text.strip())
                    if len('\n'.join(content)) > self.max_content_length:
                        break
            
            return '\n'.join(content)
        except Exception as e:
            logger.error(f"提取Word文档内容失败 {file_path}: {str(e)}")
            return None
    
    def extract_pdf_content(self, file_path: Path) -> Optional[str]:
        """提取PDF文档内容"""
        try:
            if not PDF_AVAILABLE:
                logger.warning("PDF处理库未安装，无法处理 .pdf 文件")
                return None
            
            content = []
            
            # 首先尝试使用 pdfplumber（更好的文本提取）
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages[:5]:  # 只处理前5页
                        text = page.extract_text()
                        if text:
                            content.append(text)
                            if len('\n'.join(content)) > self.max_content_length:
                                break
            except Exception:
                # 如果 pdfplumber 失败，使用 PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for i, page in enumerate(reader.pages[:5]):  # 只处理前5页
                        text = page.extract_text()
                        if text:
                            content.append(text)
                            if len('\n'.join(content)) > self.max_content_length:
                                break
            
            return '\n'.join(content)
        except Exception as e:
            logger.error(f"提取PDF文档内容失败 {file_path}: {str(e)}")
            return None
    
    def extract_excel_content(self, file_path: Path) -> Optional[str]:
        """提取Excel文档内容"""
        try:
            if not EXCEL_AVAILABLE:
                logger.warning("openpyxl 未安装，无法处理 Excel 文件")
                return None
            
            workbook = load_workbook(file_path, data_only=True)
            content = []
            
            for sheet_name in workbook.sheetnames[:3]:  # 只处理前3个工作表
                sheet = workbook[sheet_name]
                content.append(f"工作表: {sheet_name}")
                
                for row in sheet.iter_rows(max_row=50, values_only=True):  # 只处理前50行
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        content.append(' | '.join(row_text))
                    
                    if len('\n'.join(content)) > self.max_content_length:
                        break
                
                if len('\n'.join(content)) > self.max_content_length:
                    break
            
            return '\n'.join(content)
        except Exception as e:
            logger.error(f"提取Excel文档内容失败 {file_path}: {str(e)}")
            return None
    
    def extract_pptx_content(self, file_path: Path) -> Optional[str]:
        """提取PowerPoint文档内容"""
        try:
            if not PPTX_AVAILABLE:
                logger.warning("python-pptx 未安装，无法处理 .pptx 文件")
                return None
            
            presentation = Presentation(file_path)
            content = []
            
            for i, slide in enumerate(presentation.slides[:10]):  # 只处理前10张幻灯片
                content.append(f"幻灯片 {i + 1}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(shape.text.strip())
                
                if len('\n'.join(content)) > self.max_content_length:
                    break
            
            return '\n'.join(content)
        except Exception as e:
            logger.error(f"提取PowerPoint文档内容失败 {file_path}: {str(e)}")
            return None
    
    def extract_content(self, file_path: Path) -> Dict[str, Any]:
        """
        提取文件内容的主方法
        
        Returns:
            包含提取结果的字典
        """
        result = {
            'file_path': str(file_path),
            'file_name': file_path.name,
            'file_type': self.get_file_type(file_path),
            'success': False,
            'content': None,
            'error': None,
            'file_size': 0
        }
        
        try:
            if not file_path.exists():
                result['error'] = "文件不存在"
                return result
            
            result['file_size'] = file_path.stat().st_size
            
            if not self.is_supported_file(file_path):
                result['error'] = f"不支持的文件类型: {file_path.suffix}"
                return result
            
            file_type = self.get_file_type(file_path)
            content = None
            
            if file_type in ['text', 'code']:
                content = self.extract_text_content(file_path)
            elif file_type == 'markdown':
                content = self.extract_markdown_content(file_path)
            elif file_type == 'document':
                if file_path.suffix.lower() == '.docx':
                    content = self.extract_docx_content(file_path)
                elif file_path.suffix.lower() == '.pdf':
                    content = self.extract_pdf_content(file_path)
            elif file_type == 'spreadsheet':
                content = self.extract_excel_content(file_path)
            elif file_type == 'presentation':
                content = self.extract_pptx_content(file_path)
            
            if content:
                result['success'] = True
                result['content'] = content[:self.max_content_length]  # 限制长度
            else:
                result['error'] = "无法提取文件内容"
            
        except Exception as e:
            result['error'] = f"提取文件内容时发生错误: {str(e)}"
            logger.error(f"文件内容提取失败 {file_path}: {str(e)}")
        
        return result 